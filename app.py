import streamlit as st
from python-pptx import Presentation
from pptx.util import Inches, Pt
from presentobot_data import presentation_data as data

def add_title_slide(prs, topic, sub_topic):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = sub_topic
    subtitle.text = f"A {topic} Presentation"

    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

def add_intro_slide(prs, topic, sub_topic):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.shapes.placeholders[1]

    title.text = "Introduction"
    content.text = data[topic]["Intro"][sub_topic]

    content.text_frame.paragraphs[0].font.size = Pt(18)

def create_content_slides(prs, topic, sub_topic):
    if "Content" in data[topic] and sub_topic in data[topic]["Content"]:
        for section in data[topic]["Content"][sub_topic]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.shapes.placeholders[1]

            title.text = section[0]
            tf = content.text_frame
            tf.text = ""

            for item in section[1]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.space_after = Pt(12)
                p.font.size = Pt(18)

def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Thank You!"
    subtitle.text = "Created with PresentoBot"

    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

def generate_presentation(topic, sub_topic):
    prs = Presentation()
    add_title_slide(prs, topic, sub_topic)
    add_intro_slide(prs, topic, sub_topic)
    create_content_slides(prs, topic, sub_topic)
    add_closing_slide(prs)
    return prs

# ----------------- STREAMLIT UI -----------------

st.set_page_config(page_title="PresentoBot", layout="centered")
st.title("🤖 PresentoBot")
st.subheader("Auto-generate PowerPoint presentations")

# Step 1: Topic selection
topic_list = list(data.keys())
selected_topic = st.selectbox("Choose a topic", topic_list)

# Step 2: Sub-topic selection
if selected_topic:
    sub_topics = data[selected_topic]["Sub Topics"]
    selected_sub_topic = st.selectbox("Choose a sub-topic", sub_topics)

    if st.button("Generate Presentation"):
        prs = generate_presentation(selected_topic, selected_sub_topic)
        file_name = f"{selected_topic}_{selected_sub_topic.replace(' ', '_')}.pptx"

        # Save to in-memory buffer
        from io import BytesIO
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        st.success(f"Presentation generated: {file_name}")
        st.download_button(
            label="📥 Download PPT",
            data=pptx_io,
            file_name=file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
